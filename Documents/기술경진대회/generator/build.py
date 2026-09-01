# -*- coding: utf-8 -*-
"""피드포워드 오토포커스 설계 검증 리포트 → 연구소 기술경진대회 발표자료
   (제1회 연구소기술경진대회_민수개발팀 윤경현_251030 양식 사용)"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK
from deckkit import *

SRC = '/tmp/claude-0/-home-user-SEQApp/c46c0310-bc87-528b-9dd4-74e4ac81297d/scratchpad/template.pptx'
OUT = sys.argv[1] if len(sys.argv) > 1 else '/tmp/claude-0/-home-user-SEQApp/c46c0310-bc87-528b-9dd4-74e4ac81297d/scratchpad/out.pptx'

prs = Presentation(SRC)
LAY = prs.slide_layouts[1]           # '3_제목만' — 좌측 네이비 사이드바 + 로고

# ---------------------------------------------------------------- 슬라이드 정리
def drop_slides(prs, idxs):
    lst = prs.slides._sldIdLst
    els = list(lst)
    for i in sorted(idxs, reverse=True):
        prs.part.drop_rel(els[i].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'))
        lst.remove(els[i])

drop_slides(prs, range(1, 20))       # 2~20번 본문 삭제, 표지/감사 슬라이드 유지
title_slide, thanks_slide = prs.slides[0], prs.slides[1]

def new(title, size=26):
    s = prs.slides.add_slide(LAY)
    title_bar(s, title, size)
    return s

# ================================================================ 1. 표지
def retext(shape, pidx, text, size=None, bold=None, font=None, color=None):
    p = shape.text_frame.paragraphs[pidx]
    if not p.runs:
        p.add_run().text = text
    else:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
    r = p.runs[0]
    if size: r.font.size = Pt(size)
    if bold is not None: r.font.bold = bold
    if font:
        r.font.name = font
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font)
    if color: r.font.color.rgb = color
    return p

for sh in title_slide.shapes:
    if sh.name == '직사각형 8':
        retext(sh, 0, '피드포워드 오토포커스', 54, True, TITLE_FONT, NAVY)
        retext(sh, 2, '민수개발팀   윤경현', 22, True, BODY_FONT, INK)
        retext(sh, 3, '2026. 09. 01', 30, True)
    elif sh.name == 'Text 1':
        sh.left, sh.width = Inches(1.55), Inches(10.3)
        p = retext(sh, 0, '선행 Confocal 측정 기반 실시간 초점 제어 — 지연 상수 설계 검증 및 D.O.F. 오차 예산 분석',
                   16, None, BODY_FONT, INK)
        p.alignment = PP_ALIGN.CENTER

# ================================================================ 2. 개발 배경 및 개발 목적
s = new('개발 배경 및 개발 목적')
kicker(s, '동축(Co-axial) AFM 의 두 가지 구조적 한계를 선행 측정 피드포워드로 해소한다')
LX, LW, RX, RW = 0.95, 5.45, 6.95, 5.65
pill(s, LX, 2.12, 2.95, 0.50, '기술 개발 필요성')
pill(s, RX, 2.12, 2.95, 0.50, '핵심 개선 방향')
LB = [
  'AFM 이 촬상 광축과 동축 — 측정 지점과 촬상 지점이 같다',
  '감지 시점에 그 지점은 이미 촬상 중 → 보정이 항상 사후(事後)',
  '루프 지연 약 2 ms · 80 mm/s 에서 160 µm 를 지난 뒤 도달',
  '추종 범위 20 µm · 기울기 20 µm/mm — 단차에서 Defocus',
  '자기 참조 폐루프라 게인을 올리면 헌팅, 낮추면 추종 지연',
]
RB = [
  'Confocal 변위센서를 촬상점 72.5 mm 전방에 선행 배치',
  '측정값에 적용 시각을 부여해 906.25 ms 지연 큐에 보관',
  '저장된 단차 데이터로 Z축을 미리 이동 — 피드포워드 개루프',
  '대역폭 문제를 타이밍 정확도 문제로 전환',
  '측정·보정 지점 분리 → 헌팅이 구조적으로 없음',
]
def bullets(slide, x, y, w, items, size=14.5, gap=9):
    return para_block(slide, x, y, w, 3.4,
        [dict(text=t, size=size, bullet='•', line_pct=125, space_after=gap) for t in items])
bullets(s, LX, 2.88, LW, LB, size=13.5, gap=11)
bullets(s, RX, 2.88, RW, RB, size=13.5, gap=11)
banner(s, 5.92, '기존 구조는 오차를 「본 뒤에」 고쳤고, 신규 구조는 오차를 「보기 전에」 고친다', 14)
footnote(s, '출처 — 「Feedforward Autofocus System 개발」 한국전광(주) 민수개발팀, 2026-07-02 · 정량 해석은 확정 조건에서 역산한 값', 6.60)

# ================================================================ 3. 핵심 성과 지표
s = new('핵심 성과 지표')
CW3, GAP = 3.65, 0.45
STATS = [
    ('예비 시간 확보', '906.25 ms', ['72.5 mm ÷ 80 mm/s', 'Z축 이동시간 0.84 ms 의 1,080배']),
    ('추종 한계', '23배 향상', ['0.43 → 10.0 µm/ms', '지연 보상 4종 + 반주기 선보정']),
    ('촬상 시점 초점 오차', 'D.O.F. 29 %', ['1.01 µm / D.O.F. 3.5 µm', '오차의 98 %가 데드밴드']),
]
for i, (head, val, lines) in enumerate(STATS):
    x = 0.95 + i * (CW3 + GAP)
    rrect(s, x, 2.18, CW3, 2.02, CARD, INK, 0.75, 0.03)
    para_block(s, x + 0.28, 2.42, CW3 - 0.5, 0.40,
               [dict(text=head, size=17, bold=True, color=INK)])
    para_block(s, x + 0.28, 2.95, CW3 - 0.5, 0.55,
               [dict(text=val, size=27, bold=True, color=NAVY, font=TITLE_FONT)])
    para_block(s, x + 0.28, 3.58, CW3 - 0.5, 0.60,
               [dict(text=t, size=12.5, line_pct=130) for t in lines])
SUB = [
    ('제어 주기 500 µs · 초과 0 회', '38,007주기 연속 측정 · 1주기 작업시간은 예산의 42 %'),
    ('실측 여유 912 배', '추종 대상 0.011 µm/ms vs 지배 상한 10 µm/ms'),
    ('추종 범위 200 µm 목표', '기존 AFM 20 µm 의 10배 · 센서 범위 2 mm 로 여유 확보'),
]
for i, (head, sub) in enumerate(SUB):
    x = 0.95 + i * (CW3 + GAP)
    rrect(s, x, 4.42, CW3, 1.06, CARD_B, None, radius=0.05)
    para_block(s, x + 0.24, 4.62, CW3 - 0.45, 0.32,
               [dict(text=head, size=14, bold=True, color=ACCENT_D)])
    para_block(s, x + 0.24, 4.99, CW3 - 0.45, 0.42,
               [dict(text=sub, size=11.5, line_pct=120)])
hline(s, 0.95, 5.82, CW)
para_block(s, 0.95, 6.02, CW, 0.75,
    [dict(text='동축 피드백의 사후 보정 한계를 선행 측정 피드포워드로 해소 — 지연 상수의 유도, 코드 반영, '
                '실측 로그가 모두 일치하며 현 요구 조건을 여유 있게 충족한다.', size=14, line_pct=130)])

# ================================================================ 4. 기존 동축 방식의 구조적 한계
s = new('기존 동축 방식의 구조적 한계')
kicker(s, '측정 지점과 촬상 지점이 같으면 미리 움직일 시간이 없다 — 세 한계는 모두 이 하나에서 나온다')
LIM = [
    ('①', '지연 (Latency)', [
        't loop = 센서·연산 1 ms + Z축 이동 0.85–1.32 ms ≈ 2 ms',
        '80 mm/s 에서 160 µm 를 지난 뒤 보정이 도달',
        'D.O.F. 소진 지점 = 3.5 µm ÷ 2 ms = 1.75 µm/ms']),
    ('②', '좁은 추적 범위', [
        'Z축 추적 가능 높이 변화 20 µm — D.O.F. 의 5.7배',
        '기울기 20 µm/mm → 스캔 80 mm/s 에서 Z축 1.6 mm/s',
        '범위 이탈 시 추적 상실 → 재포착까지 검사 불능']),
    ('③', '급격한 단차', [
        '자재 딥 100–160 µm = D.O.F. 3.5 µm 의 30–45배',
        '딥 추종에 주행 8 mm · 시간 100 ms 필요 (기존 AFM)',
        '추종하지 못하고 Defocus 발생']),
]
for i, (num, head, lines) in enumerate(LIM):
    x = 0.95 + i * (CW3 + GAP)
    rrect(s, x, 2.18, CW3, 2.95, CARD, INK, 0.75, 0.03)
    badge = rrect(s, x + 0.26, 2.42, 0.52, 0.52, ACCENT_D, None, shape=MSO_SHAPE.OVAL)
    label_in(badge, num, 17, True, WHITE)
    para_block(s, x + 0.92, 2.53, CW3 - 1.15, 0.36,
               [dict(text=head, size=16, bold=True, color=NAVY)])
    para_block(s, x + 0.26, 3.12, CW3 - 0.52, 1.75,
               [dict(text=t, size=12, bullet='·', line_pct=120, space_after=7) for t in lines])
warn = rrect(s, 0.95, 5.32, CW, 0.86, RGBColor(0xFB, 0xEF, 0xEC), RED, 0.75, 0.05)
label_in(warn, '이 하한은 소프트웨어로 개선되지 않는다 — 줄이려면 더 빠른 센서와 Z축이 필요하다. '
               '루프를 빠르게 하는 대신 구조를 바꾸는 것이 유일한 해법이다.', 14, True, RED, PP_ALIGN.CENTER)
footnote(s, '완만한 휨(요구 1.0 µm/mm)은 기존 AFM 으로도 담기는 영역이었다 — 실제 문제는 딥·급단차의 국부 기울기와 사후 보정 지연이다.', 6.46)

# ================================================================ 5. 피드포워드 전환 — 선행 배치
s = new('피드포워드 전환 — 선행 측정 구조')
sen = rrect(s, 2.05, 2.20, 2.75, 1.00, ACCENT_D, None, radius=0.06)
label_in(sen, 'Confocal 센서 FS2404-2\n측정 범위 2 mm', 13.5, True, WHITE, line_pct=125)
cam = rrect(s, 8.55, 2.20, 2.75, 1.00, NAVY, None, radius=0.06)
label_in(cam, 'Line-scan 카메라\nD.O.F. 3.5 µm', 13.5, True, WHITE, line_pct=125)
para_block(s, 4.90, 2.28, 3.55, 0.34,
           [dict(text='72.5 mm', size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=TITLE_FONT)])
arrow(s, 4.95, 2.72, 3.45, 0.30, ACCENT, MSO_SHAPE.LEFT_RIGHT_ARROW)
bar = rect(s, 1.35, 3.52, 10.55, 0.62, RGBColor(0xE9, 0xEC, 0xF1), INK, 0.75)
label_in(bar, '자재 (패널 · 전체 길이 약 200 mm · 높이 편차 목표 200 µm)', 13, False, INK)
for cx in (3.425, 9.925):
    rect(s, cx - 0.01, 3.20, 0.022, 0.33, ACCENT_D)
sc = arrow(s, 1.35, 4.30, 3.05, 0.34, ACCENT_D)
label_in(sc, '스캔 방향', 12, True, WHITE)
para_block(s, 4.60, 4.34, 5.0, 0.30, [dict(text='X축 80 mm/s (연속 이동)', size=13, bold=True, color=INK)])
fb = rrect(s, 0.95, 4.94, CW, 0.86, ACCENT_L, None, radius=0.05)
label_in(fb, '906.25 ms  =  72.5 mm ÷ 80 mm/s     —     센서가 읽은 지점이 촬상점에 도착하기까지의 예비 시간',
         17, True, NAVY, PP_ALIGN.CENTER, font=TITLE_FONT)
NOTE = [
    ('Z축 이동시간의 1,080배', '축이 목표에 도달하고도 한참을 기다린 뒤에 그 지점이 촬상된다. 대역폭이 아니라 타이밍 정확도가 성능을 정한다.'),
    ('헌팅이 구조적으로 없다', '측정과 보정이 다른 지점에서 일어나므로 Z를 움직여도 센서 값이 흔들리지 않는다 — 자기 참조 폐루프가 사라진다.'),
]
for i, (h, t) in enumerate(NOTE):
    x = 0.95 + i * 6.05
    rrect(s, x, 5.98, 5.85, 0.94, CARD, None, radius=0.05)
    para_block(s, x + 0.25, 6.14, 5.35, 0.28, [dict(text=h, size=13.5, bold=True, color=ACCENT_D)])
    para_block(s, x + 0.25, 6.46, 5.35, 0.42, [dict(text=t, size=11.5, line_pct=120)])

# ================================================================ 6. 시스템 구성
s = new('시스템 구성 — 측정 경로와 지령 경로')
def chain(slide, y, items, x0=0.95, total=None, h=0.92, fill=ACCENT_D, tc=WHITE, size=12):
    total = CW if total is None else total
    n = len(items)
    ov = 0.16
    w = (total + ov * (n - 1)) / n
    for i, (a, b) in enumerate(items):
        sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x0 + i * (w - ov)),
                                    Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        no_line(sh); sh.shadow.inherit = False
        label_in(sh, a + '\n' + b, size, True, tc, line_pct=120)
        sh.text_frame.paragraphs[1].runs[0].font.size = Pt(size - 2)
        sh.text_frame.paragraphs[1].runs[0].font.bold = False
para_block(s, 0.95, 2.02, 6.0, 0.28, [dict(text='측정 경로 — 아날로그', size=13.5, bold=True, color=ACCENT_D)])
chain(s, 2.34, [('자재 표면', '높이 편차 200 µm'), ('FS2404-2', 'Confocal 프로브'),
                ('IFC2421', '0–10 V = 0–2 mm'), ('SIO-AI8F', '8채널 아날로그 입력'),
                ('제어 PC', 'AxaiSwReadVoltage')], fill=ACCENT_D)
mid = rrect(s, 0.95, 3.46, CW, 1.02, ACCENT_L, ACCENT_D, 0.75, 0.04)
label_in(mid, '제어 PC — HiPrecisionLoop 500 µs (QPC 기준)\n'
              'IIR 필터 α = 0.10  →  데드밴드 ±5 mV  →  목표 위치 환산 ×2,000 pulse/V     ·     지연 큐 906.25 ms · 대기 약 1,813개',
         13, True, NAVY, line_pct=125)
mid.text_frame.paragraphs[1].runs[0].font.size = Pt(12)
para_block(s, 0.95, 4.64, 6.0, 0.28, [dict(text='지령 경로 — 펄스열', size=13.5, bold=True, color=NAVY)])
chain(s, 4.96, [('AXT SMC-2V04', '펄스 출력 최대 500 kpps'), ('Ezi-STEP MI 28L', 'FASTECH 스테퍼 드라이브'),
                ('BM 28L · Z축', '1 pulse = 0.1 µm'), ('카메라 초점', 'D.O.F. 예산 2.5 µm')], fill=NAVY)
para_block(s, 0.95, 6.14, CW, 0.75, [
    dict(text='회로도(KEOC2604) ↔ 실행 시 소프트웨어 열거 대조 — 축 2 / AI 8채널 / AXT_PCIeB_LBExR : 전 항목 일치',
         size=12.5, bold=True, color=ACCENT_D, space_after=4),
    dict(text='엔코더 되먹임 배선 없음 — AxmStatusGetActPos 는 보드가 내보낸 펄스 누적값이다. 개루프라 헌팅은 없지만 Z축 실위치는 컨트롤러에서 검증할 수 없다.',
         size=12, color=INK_SOFT, line_pct=120)])

# ================================================================ 7. 지연 상수 유도
s = new('지연 상수 906.25 ms 의 유도')
fx = rrect(s, 0.95, 2.14, 5.55, 1.10, CARD, ACCENT_D, 0.75, 0.04)
label_in(fx, 't delay  =  센서–카메라 거리 ÷ X축 속도\n= 72.5 mm ÷ 80 mm/s  =  906.25 ms',
         15, True, NAVY, line_pct=135, font=TITLE_FONT)
bullets(s, 0.95, 3.40, 5.55, [
    'applyTime = sampleTime + 906.25 ms — 절대 시각 기준',
    '주기를 1 ms → 500 µs 로 바꿔도 지연은 그대로 유지',
    '「906 샘플」은 1 ms 주기 표현 · 500 µs 에서는 1,812 샘플',
    'Setting.ini FocusDelayMs = 906.250 (double)',
], size=12.5, gap=8)
qx, qw = 6.95, 5.83
para_block(s, qx, 2.14, qw, 0.28, [dict(text='지연 큐 — 시간 기준 선형 버퍼', size=13.5, bold=True, color=ACCENT_D)])
rect(s, qx, 2.52, 1.55, 0.72, RGBColor(0xDD, 0xDD, 0xDD), INK, 0.5)
rect(s, qx + 1.55, 2.52, 2.75, 0.72, ACCENT, INK, 0.5)
rect(s, qx + 4.30, 2.52, qw - 4.30, 0.72, WHITE, INK, 0.5)
para_block(s, qx, 2.72, 1.55, 0.3, [dict(text='소비 완료', size=10.5, align=PP_ALIGN.CENTER)])
para_block(s, qx + 1.55, 2.72, 2.75, 0.3, [dict(text='대기 중 · 약 1,813개', size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)])
para_block(s, qx + 4.30, 2.72, qw - 4.30, 0.3, [dict(text='빈 공간', size=10.5, color=INK_SOFT, align=PP_ALIGN.CENTER)])
para_block(s, qx + 0.85, 3.30, 1.6, 0.5, [dict(text='readIndex\n906 ms 전 값', size=10, color=ACCENT_D, align=PP_ALIGN.CENTER, line_pct=115)])
para_block(s, qx + 3.45, 3.30, 1.7, 0.5, [dict(text='writeIndex\n지금 측정한 값', size=10, color=ACCENT_D, align=PP_ALIGN.CENTER, line_pct=115)])
para_block(s, qx, 3.95, qw, 0.55, [dict(text='두 인덱스의 간격이 곧 지연 상수다. 링버퍼가 아니라 한 방향으로 전진하는 선형 버퍼(100,000 슬롯 = 50초분)이며, 끝에 닿으면 미소비 구간을 앞으로 압축한다.', size=11.5, line_pct=125)])
eff = rrect(s, 0.95, 5.26, CW, 0.74, ACCENT_L, None, radius=0.05)
label_in(eff, '부수 효과 — 주기를 1 ms → 500 µs 로 줄이면서 X 방향 높이 샘플 간격이 80 µm → 40 µm 로 조밀해졌다 (짧은 주기의 요철 포착)',
         13.5, True, NAVY, PP_ALIGN.CENTER)
para_block(s, 0.95, 6.12, CW, 0.86, [
    dict(text='촬상 조건 (정상 획득 이미지 3,200 × 600,000 px · 8-bit · 1.8 GB 에서 역산)', size=12.5, bold=True, color=ACCENT_D, space_after=5),
    dict(text='스캔 라인 피치 0.333 µm/라인   ·   라인 레이트 240 kHz   ·   1회 스캔 2.5 s   ·   초점 갱신 간격 120 라인(X축 40 µm)   ·   선행 구간 약 217,500 라인',
         size=12, line_pct=125)])

# ================================================================ 8. 제어 주기 실측 검증
s = new('제어 주기 실측 검증 — 500 µs')
kicker(s, '실장비 측정값 · 실앱 20초 · 카메라 2채널 동시 구동 · 매 주기 로깅 OFF 조건')
T61 = [('평균 주기', '500 µs', '40개 구간 전부 정확히 500'),
       ('주기 초과', '0 회 / 38,007주기', '오차가 누적되지 않는 절대 시각 만기'),
       ('주기 지터', '대부분 ±15 µs', '최악 구간 340–660 µs'),
       ('1주기 작업시간', '평균 39.8 µs', 'p99 45.2 µs / 최대 211 µs'),
       ('예산 사용률', '42 %', '최대치 211 µs ÷ 500 µs')]
tb = table(s, 0.95, 2.30, 6.35, 2.75, 6, 3, col_w=[2.0, 2.0, 3.2], header_h=0.44, row_h=0.44)
for c, h in enumerate(['항목', '결과', '조건']):
    cell(tb, 0, c, h, 12.5, True, WHITE, ACCENT_D, PP_ALIGN.LEFT)
for r, row in enumerate(T61, 1):
    for c, v in enumerate(row):
        cell(tb, r, c, v, 11.5, c == 1, NAVY if c == 1 else INK,
             RGBColor(0xF7, 0xF9, 0xFC) if r % 2 else None)
style_table(tb)
cd = CategoryChartData()
cd.categories = ['AxmSignalReadInput', 'AxmStatusGetActPos', 'AxmStatusReadMotion', 'AxaiSwReadVoltage']
cd.add_series('평균 소요 시간 (µs)', (17.2, 9.4, 9.3, 3.8))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(7.62), Inches(2.20),
                        Inches(5.16), Inches(2.95), cd)
ch = gf.chart
ch.has_title = True
ch.chart_title.text_frame.text = '주기당 AXL API 소요 시간 (µs) · 합계 39.7'
set_para(ch.chart_title.text_frame.paragraphs[0], 12, True, NAVY, PP_ALIGN.CENTER)
ch.has_legend = False
pl = ch.plots[0]; pl.gap_width = 60
pl.has_data_labels = True
dl = pl.data_labels; dl.font.size = Pt(10.5); dl.font.bold = True; dl.font.color.rgb = INK
dl.position = XL_LABEL_POSITION.OUTSIDE_END
ser = pl.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb = ACCENT
ca, va = ch.category_axis, ch.value_axis
ca.tick_labels.font.size = Pt(10.5); ca.tick_labels.font.color.rgb = INK
ca.has_major_gridlines = False; ca.major_tick_mark = XL_TICK_MARK.NONE
va.tick_labels.font.size = Pt(10); va.tick_labels.font.color.rgb = INK_SOFT
va.has_major_gridlines = True
va.format.line.color.rgb = GREY_L
va.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xEC, 0xEC)
para_block(s, 0.95, 5.36, CW, 0.95, [
    dict(text='개별 주기가 튀어도 평균이 밀리지 않는 것은 만기를 절대 시각으로 잡기 때문이다 — 한 주기가 늦으면 다음 주기가 그만큼 짧아져 오차가 누적되지 않는다.',
         size=13.5, line_pct=130, space_after=7),
    dict(text='미검증 — 모션 명령(AxmOverridePosAbs / AxmMoveStartPos)의 소요 시간은 축이 실제로 움직여야 측정 가능해 아직 재지 않았다. 수락·거부율은 실측 로그로 정량화했다.',
         size=12, color=INK_SOFT, line_pct=125)])

# ================================================================ 9. 지연 오차 예산
s = new('지연 오차 예산 — 성분과 보상 상태')
kicker(s, 'e focus = Δt(지연 오차) × dz/dt(표면 높이 변화율) — 예측되고 보상되는 시간은 오차가 아니다')
T71 = [('IIR 필터 군지연  (1−α)/α = 9샘플', '4.50 ms', '적용 시각에서 filterDelayUs 만큼 차감', '보상됨'),
       ('틱 양자화  GetTickCount64 1 ms', '±1.00 ms', '큐 시각을 QPC 기준 applyTimeUs 로 전환', '제거됨'),
       ('지연 상수 절삭  906.25 → 906', '0.25 ms', 'delayTimeMs 를 double 로, 기본값 906.25', '제거됨'),
       ('Lead 올림 절삭  기존 ceil() 정수 ms', '최대 1.00 ms', 'µs 단위 연산으로 전환', '제거됨'),
       ('Z축 이동 시간  35 pulse · a = 200 M', '0.84 ms', 'Lead Control 로 이동시간만큼 선행 지령', '보상됨'),
       ('만기 판정 편향  항상 늦는 쪽으로만', '+0.25 ms', '적용 시각에서 반주기를 미리 차감', '제거됨'),
       ('제어 주기 양자화  만기 판정 500 µs 격자', '±0.25 ms', '주기를 줄이는 것 외에 방법 없음', '잔존')]
tb = table(s, 0.95, 2.28, CW, 3.30, 8, 4, col_w=[4.1, 1.5, 4.6, 1.1], header_h=0.42, row_h=0.41)
for c, h in enumerate(['지연 성분', '크기', '처리', '상태']):
    cell(tb, 0, c, h, 12.5, True, WHITE, ACCENT_D, PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
for r, row in enumerate(T71, 1):
    resid = row[3] == '잔존'
    for c, v in enumerate(row):
        cell(tb, r, c, v, 11.5, c in (1, 3), 
             (ORANGE if resid else RGBColor(0x1F, 0x6B, 0x4A)) if c == 3 else (NAVY if c == 1 else INK),
             RGBColor(0xFD, 0xF3, 0xE8) if resid else (RGBColor(0xF7, 0xF9, 0xFC) if r % 2 else None),
             PP_ALIGN.CENTER if c in (1, 3) else PP_ALIGN.LEFT)
style_table(tb)
a = rrect(s, 0.95, 5.72, 6.35, 1.05, CARD, None, radius=0.05)
label_in(a, '적용 시각 계산\n906,250 − 4,500(군지연) − 840(Z축 이동) − 250(반주기)  =  900,660 µs',
         13, True, NAVY, PP_ALIGN.CENTER, line_pct=130)
a.text_frame.paragraphs[0].runs[0].font.size = Pt(11.5)
a.text_frame.paragraphs[0].runs[0].font.color.rgb = ACCENT_D
b = rrect(s, 7.62, 5.72, 5.16, 1.05, RGBColor(0xFB, 0xEF, 0xEC), RED, 0.75, 0.05)
label_in(b, '주의 — 상수를 손으로 빼지 말 것\nFocusDelayMs 에는 물리값 906.250 을 그대로 넣는다. 보상은 코드가 한다.',
         12, True, RED, PP_ALIGN.CENTER, line_pct=130)
b.text_frame.paragraphs[1].runs[0].font.size = Pt(11)
b.text_frame.paragraphs[1].runs[0].font.color.rgb = INK

# ================================================================ 10. 반주기 선보정
s = new('반주기 선보정 — 코드 한 줄로 Δt 절반')
code = rrect(s, 0.95, 2.14, CW, 0.60, RGBColor(0x2C, 0x31, 0x4F), None, radius=0.05)
label_in(code, 'offsetUs  -=  FOCUS_LOOP_PERIOD_US / 2.0;', 15, True, WHITE, PP_ALIGN.CENTER, font='Courier New')
BA = [(0.95, 5.05, '보정 전', RGBColor(0xF3, 0xF1, 0xF1), INK,
       ['적용 오차 ∈ [ 0, +500 ) µs', '평균 +250 µs · 최대 500 µs', '늦는 쪽으로만 쌓이는 계통 편향']),
      (7.05, 5.73, '보정 후 — 현재 구현', ACCENT_L, NAVY,
       ['적용 오차 ∈ ( −250, +250 ] µs', '평균 0 µs · 최대 250 µs · RMS 144 µs', '부호가 아니라 크기로 소모되므로 유리'])]
for x, w, head, fill, hc, lines in BA:
    rrect(s, x, 2.96, w, 1.62, fill, None, radius=0.04)
    para_block(s, x + 0.28, 3.16, w - 0.55, 0.32, [dict(text=head, size=15, bold=True, color=hc)])
    para_block(s, x + 0.28, 3.62, w - 0.55, 0.85,
               [dict(text=t, size=12.5, bullet='·', line_pct=125, space_after=6) for t in lines])
arrow(s, 6.14, 3.55, 0.72, 0.44)
para_block(s, 0.95, 4.78, CW, 0.30,
           [dict(text='만기 판정 nowUs >= applyTimeUs 는 500 µs 격자에서만 일어나 항상 다음 주기에 걸린다 — 반주기를 미리 빼두면 평균이 0 이 되고 최대 절대오차가 절반이 된다.',
                 size=12.5, color=INK_SOFT)])
# 시간축
rect(s, 1.35, 5.72, 10.4, 0.16, ACCENT_L)
rect(s, 10.95, 5.72, 0.80, 0.16, ACCENT)
for x, lab in [(1.35, '측정 시각'), (10.90, '지령 시각'), (11.72, '촬상 시각')]:
    rect(s, x - 0.015, 5.60, 0.03, 0.40, NAVY)
para_block(s, 1.35, 5.28, 3.0, 0.28, [dict(text='측정 시각', size=12, bold=True, color=NAVY)])
para_block(s, 8.10, 5.28, 2.7, 0.28, [dict(text='지령 시각  측정 + 900,660 µs', size=12, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)])
para_block(s, 10.30, 6.06, 2.5, 0.28, [dict(text='촬상 시각  + 906,250 µs', size=12, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)])
para_block(s, 1.35, 6.06, 7.0, 0.28, [dict(text='906.25 ms 예비 시간 — 이 구간 동안 값을 보관한다', size=12, color=INK_SOFT)])
para_block(s, 0.95, 6.52, CW, 0.36, [dict(text='보상 합계 5,590 µs = 필터 군지연 4,500 + Z축 이동 840 + 루프 반주기 250 — 전체 906 ms 의 0.62 % 구간이다.', size=12, color=INK_SOFT)])

# ================================================================ 11. 추종 한계 비교
s = new('추종 가능한 최대 높이 변화율')
cd = CategoryChartData()
cd.categories = ['보상 전\nΔt = 5.75 ms', '필터만 보상\nΔt = 1.25 ms', '지연 보상 3종\nΔt = 0.50 ms',
                 '현재 구현\nΔt = ±0.25 ms', '＋데드밴드 SKIP\n예산 3.5 µm']
cd.add_series('추종 한계 (µm/ms)', (0.43, 2.00, 5.00, 10.00, 14.00))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.95), Inches(2.12),
                        Inches(7.35), Inches(4.10), cd)
ch = gf.chart
ch.has_title = True
ch.chart_title.text_frame.text = '초점 정확도 기준 · D.O.F. 예산 2.5 µm ÷ 잔존 시간 오차 (µm/ms)'
set_para(ch.chart_title.text_frame.paragraphs[0], 12, True, NAVY, PP_ALIGN.CENTER)
ch.has_legend = False
pl = ch.plots[0]; pl.gap_width = 55; pl.has_data_labels = True
dl = pl.data_labels; dl.font.size = Pt(11); dl.font.bold = True; dl.font.color.rgb = INK
dl.position = XL_LABEL_POSITION.OUTSIDE_END
ser = pl.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb = RGBColor(0xC5, 0xD6, 0xE8)
for i, col in {3: ACCENT_D, 4: RGBColor(0xA8, 0xC2, 0xDC)}.items():
    pt = ser.points[i]; pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
ca, va = ch.category_axis, ch.value_axis
ca.tick_labels.font.size = Pt(10.5); ca.tick_labels.font.color.rgb = INK
ca.has_major_gridlines = False; ca.major_tick_mark = XL_TICK_MARK.NONE
va.tick_labels.font.size = Pt(10); va.tick_labels.font.color.rgb = INK_SOFT
va.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xEC, 0xEC)
va.format.line.color.rgb = GREY_L
RX2, RW2 = 8.58, 4.20
para_block(s, RX2, 2.12, RW2, 0.30, [dict(text='네 가지 상한 — 가장 작은 값이 실제 한계', size=13, bold=True, color=ACCENT_D)])
UP = [('① Z축 명령 상한', '50 µm/ms', 'FASTECH 500,000 pulse/s · 하드웨어 절대 한계', False),
      ('② Z축 속도 설정', '15 µm/ms', 'Vel 150,000 pulse/s · 설정으로 조정 가능', False),
      ('②′ Z축 가속도', '10 µm/ms', 'Accel 200 M · 과도 응답 상한', True),
      ('③ 초점 정확도', '10.0 µm/ms', 'D.O.F. 예산 2.5 µm ÷ 0.25 ms · 정상상태', True)]
yy = 2.46
for head, val, sub, hot in UP:
    rrect(s, RX2, yy, RW2, 0.82, ACCENT_L if hot else CARD, ACCENT_D if hot else None, 0.75, 0.05)
    para_block(s, RX2 + 0.22, yy + 0.13, RW2 - 0.44, 0.28, [dict(text=head, size=12, bold=True, color=INK)])
    para_block(s, RX2 + 0.22, yy + 0.13, RW2 - 0.44, 0.28, [dict(text=val, size=13, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)])
    para_block(s, RX2 + 0.22, yy + 0.46, RW2 - 0.44, 0.30, [dict(text=sub, size=10.5, color=INK_SOFT, line_pct=115)])
    yy += 0.90
tag = rrect(s, RX2, yy + 0.12, RW2, 0.52, NAVY, None, radius=0.06)
label_in(tag, '지배 상한  10 µm/ms', 14, True, WHITE)
para_block(s, 0.95, 6.34, 7.35, 0.60,
    [dict(text='지배 상한에서 ②′와 ③이 같은 값에서 만난다 — 구동 능력과 초점 정확도가 동시에 한계에 닿는 균형점이라 '
                '어느 한쪽만 올려서는 전체가 오르지 않는다.', size=11, color=INK_SOFT, line_pct=120)])

# ================================================================ 12. 오차 예산 vs D.O.F.
s = new('오차 예산과 D.O.F. 판정')
T73 = [('데드밴드 (USE 시 고정)', '1.00 µm', '1.00 µm', '1.00 µm'),
       ('제어 주기 양자화 ±0.25 ms', '0.01 µm', '0.02 µm', '0.25 µm'),
       ('필터 · 절삭 · 틱 · 편향', '0 µm', '0 µm', '0 µm'),
       ('합계', '1.01 µm', '1.02 µm', '1.25 µm'),
       ('데드밴드 SKIP 시 합계', '0.01 µm', '0.02 µm', '0.25 µm'),
       ('D.O.F. 3.5 µm 대비', '29 % 통과', '29 % 통과', '36 % 통과')]
tb = table(s, 0.95, 2.18, 7.55, 3.05, 7, 4, col_w=[3.1, 1.5, 1.5, 1.5], header_h=0.55, row_h=0.41)
for c, h in enumerate(['성분', '실측 조건\n0.04 µm/ms', '개발 목표\n0.08 µm/ms', '가정 조건\n1.0 µm/ms']):
    cell(tb, 0, c, h, 11.5, True, WHITE, ACCENT_D, PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
for r, row in enumerate(T73, 1):
    strong = row[0] in ('합계', 'D.O.F. 3.5 µm 대비')
    for c, v in enumerate(row):
        cell(tb, r, c, v, 11.5, strong or c == 0 and False,
             NAVY if strong else INK,
             ACCENT_L if strong else (RGBColor(0xF7, 0xF9, 0xFC) if r % 2 else None),
             PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
style_table(tb)
GX, GW = 8.85, 3.93
para_block(s, GX, 2.18, GW, 0.30, [dict(text='D.O.F. 예산 사용률 · 개발 목표 조건', size=12.5, bold=True, color=ACCENT_D)])
rect(s, GX, 2.58, GW, 0.62, RGBColor(0xEE, 0xEE, 0xEE), INK, 0.75)
rect(s, GX, 2.58, GW * 1.02 / 3.5, 0.62, ACCENT_D)
para_block(s, GX, 3.28, GW, 0.28, [dict(text='1.02 µm', size=12, bold=True, color=NAVY),])
para_block(s, GX, 3.28, GW, 0.28, [dict(text='D.O.F. 3.5 µm', size=12, color=INK_SOFT, align=PP_ALIGN.RIGHT)])
big = rrect(s, GX, 3.75, GW, 1.05, ACCENT_L, None, radius=0.05)
label_in(big, '29 %', 34, True, NAVY, font=TITLE_FONT)
para_block(s, GX, 4.96, GW, 0.60, [dict(text='오차의 98 %가 데드밴드다. SKIP 하면 합계 0.02 µm — 예산의 1 % 미만이 된다.', size=12, line_pct=125)])
warn = rrect(s, 0.95, 5.52, CW, 0.72, RGBColor(0xFD, 0xF3, 0xE8), ORANGE, 0.75, 0.05)
label_in(warn, '단, 계단 응답은 별개다 — 1차 IIR 의 95 % 정착은 3τ = 14.2 ms 이고 그동안 X축은 1.14 mm(약 3,400 라인)를 지나간다. '
               '큰 단차는 유효 대역에서 기각되므로 문제 구간은 대역 안에서 빠르게 변하는 곳뿐이다.', 12, True, ORANGE, PP_ALIGN.CENTER)
footnote(s, '지연 보상 전 예측 1.24 µm 에서 현재 1.01 µm 로 내려갔고, 가정 조건 1 µm/ms 도 초과(200 %)에서 통과(36 %)로 바뀌었다 — '
            '계산이 「정상 이미지 획득」이라는 현장 결과와 일치한다.', 6.42)

# ================================================================ 13. 자재 단차 대응
s = new('자재 단차 대응과 초기 위치 정렬')
kicker(s, '앞 절이 「정상 신호일 때 얼마나 정확한가」라면, 여기는 신호가 깨질 때와 스캔이 시작될 때다')
DEAL = [(0.95, 5.70, '① 유효 대역 기각', 'VOLTAGE_RANGE_REJECT', [
            '딥 100–160 µm = D.O.F. 3.5 µm 의 30–45배 — 주변 표면과 동시에 초점을 맞출 수 없다',
            '추종하면 왕복 22.8 ms 동안 X축이 1.83 mm 를 지나가 앞뒤 정상 구간까지 손실',
            '대역 밖 값은 채택하지 않고 직전 유효 목표를 재저장 — 지연 큐의 시간축은 그대로 유지',
            '급정지가 아니라 마지막 유효 목표 유지 — 단차가 끝나는 즉시 추종 재개']),
        (7.08, 5.70, '② 첫 데이터 즉시 적용', 'First Data Immediate', [
            '패널마다 두께·안착 높이가 달라 스캔 시작 시점의 Z 오프셋이 제각각이다',
            '첫 유효 샘플을 지연 없이 즉시 적용하고, 두 번째부터 906 ms 지연 큐를 태운다',
            '미적용 시 초기 오프셋 100 µm 기준 패널 앞 0.59 mm 가 초점을 벗어난다',
            '트리거 시점 카메라는 72.5 mm 뒤 — Z 이동 7.4 ms 의 122배 여유를 그대로 쓴다'])]
for x, w, head, code, lines in DEAL:
    rrect(s, x, 2.12, w, 3.32, CARD, INK, 0.75, 0.03)
    para_block(s, x + 0.28, 2.34, w - 0.55, 0.32, [dict(text=head, size=16, bold=True, color=NAVY)])
    para_block(s, x + 0.28, 2.74, w - 0.55, 0.28, [dict(text=code, size=11.5, bold=True, color=ACCENT_D, font='Courier New')])
    para_block(s, x + 0.28, 3.12, w - 0.55, 2.10,
               [dict(text=t, size=12, bullet='·', line_pct=120, space_after=8) for t in lines])
para_block(s, 0.95, 5.58, CW, 0.30, [dict(text='유효 대역 설정 기준 — 정상 표면보다 넓고 자재 단차보다 좁게', size=13, bold=True, color=ACCENT_D)])
BAND = [('실측 신호 전 범위', '1.243 – 1.413 mm', '폭 170 µm · 정상 표면 + 자재 단차'),
        ('정상 표면대 (추정)', '약 1.37 – 1.41 mm', '폭 약 40 µm · 추종 대상'),
        ('자재 단차', '100 – 160 µm', 'D.O.F. 의 30–45배 · 대역 밖으로 기각')]
for i, (h, v, t) in enumerate(BAND):
    x = 0.95 + i * (CW3 + GAP)
    rrect(s, x, 5.92, CW3, 0.96, CARD_B, None, radius=0.05)
    para_block(s, x + 0.22, 6.08, CW3 - 0.44, 0.26, [dict(text=h, size=11.5, bold=True, color=ACCENT_D)])
    para_block(s, x + 0.22, 6.34, CW3 - 0.44, 0.26, [dict(text=v, size=13, bold=True, color=NAVY)])
    para_block(s, x + 0.22, 6.60, CW3 - 0.44, 0.24, [dict(text=t, size=10.5, color=INK_SOFT)])

# ================================================================ 14. Z축 구동 능력
s = new('Z축 구동 능력 — 병목은 구동이 아니다')
T92 = [('20 pulse', '2 µm/ms', '40,000,000', '20,000', '5배 여유'),
       ('40 pulse', '4 µm/ms', '80,000,000', '40,000', '2.5배 여유'),
       ('60 pulse', '6 µm/ms', '120,000,000', '60,000', '1.7배 여유'),
       ('100 pulse', '10 µm/ms', '200,000,000', '100,000', '현재 설정과 정확히 일치 · 지배 상한'),
       ('150 pulse', '15 µm/ms', '300,000,000', '150,000', '가속도 부족 (속도는 충족)'),
       ('200 pulse', '20 µm/ms', '400,000,000', '200,000', '둘 다 부족')]
tb = table(s, 0.95, 2.18, 7.75, 3.10, 7, 5, col_w=[1.35, 1.35, 1.85, 1.45, 3.0], header_h=0.55, row_h=0.42)
for c, h in enumerate(['1 ms 변화량', '높이 변화율', '필요 가속도\n[pulse/s²]', '필요 속도\n[pulse/s]', '현재 설정 대비\nAccel 200 M · Vel 150,000']):
    cell(tb, 0, c, h, 10.5, True, WHITE, ACCENT_D, PP_ALIGN.CENTER)
for r, row in enumerate(T92, 1):
    hot = r == 4
    for c, v in enumerate(row):
        cell(tb, r, c, v, 11, hot, NAVY if hot else INK,
             ACCENT_L if hot else (RGBColor(0xF7, 0xF9, 0xFC) if r % 2 else None),
             PP_ALIGN.LEFT if c == 4 else PP_ALIGN.CENTER)
style_table(tb)
RX3, RW3 = 8.98, 3.80
ZC = [('D.O.F. 한 칸 보정', '0.84 ms', 't = 2√(35 ÷ 2×10⁸) · 삼각 프로파일 · 최대 83,700 pulse/s = 상한의 17 %'),
      ('그동안 시료 이동', '67 µm', '제어 주기 2회 · 높이 샘플 2개 · 906 ms 예비 시간의 1/1,080'),
      ('먼저 걸리는 쪽', '가속도', 'Vel 150,000 은 15 µm/ms 까지, Accel 200 M 은 10 µm/ms 까지 감당')]
yy = 2.18
for h, v, t in ZC:
    rrect(s, RX3, yy, RW3, 1.00, CARD, None, radius=0.05)
    para_block(s, RX3 + 0.22, yy + 0.14, RW3 - 0.44, 0.26, [dict(text=h, size=11.5, bold=True, color=ACCENT_D)])
    para_block(s, RX3 + 0.22, yy + 0.14, RW3 - 0.44, 0.26, [dict(text=v, size=13, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)])
    para_block(s, RX3 + 0.22, yy + 0.48, RW3 - 0.44, 0.44, [dict(text=t, size=10.5, color=INK_SOFT, line_pct=120)])
    yy += 1.07
banner(s, 5.54, '10 µm/ms 에서 필요 가속도 200 M 과 현재 설정이 정확히 일치하고, 초점 정확도(2.5 µm ÷ 0.25 ms)도 같은 지점에서 만족된다', 13.5)
footnote(s, '다만 여유가 정확히 0 인 균형점이다 — 초점 오차 2.5 µm + 데드밴드 1.0 µm = D.O.F. 3.5 µm. 사양으로 확정하려면 데드밴드 SKIP 으로 1 µm 를 회수해 두는 것이 안전하다.', 6.30)

# ================================================================ 15. 실측 로그 검증
s = new('실측 로그 검증 — 36,472행 · 6.9 MB')
kicker(s, 'AutoFocus_20260813_144610.csv · 시퀀스 5,999 ms · 유효 대역 4.00–4.55 V · 사이클 9,340 / 유효 표본 3,025')
cd = CategoryChartData()
cd.categories = ['평탄 1', '평탄 2', '평탄 3', '평탄 4', '평탄 5']
cd.add_series('평탄부 중앙값 높이 (µm)', (847.1, 840.3, 846.2, 851.9, 858.2))
gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.95), Inches(2.08),
                        Inches(6.90), Inches(3.15), cd)
ch = gf.chart
ch.has_title = True
ch.chart_title.text_frame.text = '추종 대상 — 단차를 걷어내면 남는 완만한 기울기'
set_para(ch.chart_title.text_frame.paragraphs[0], 12, True, NAVY, PP_ALIGN.CENTER)
ch.has_legend = False
ser = ch.plots[0].series[0]
ser.format.line.color.rgb = ACCENT_D; ser.format.line.width = Pt(2.25)
ser.smooth = False
ch.plots[0].has_data_labels = True
dl = ch.plots[0].data_labels; dl.font.size = Pt(10.5); dl.font.bold = True; dl.font.color.rgb = INK
dl.position = XL_LABEL_POSITION.ABOVE
va = ch.value_axis
va.minimum_scale = 830; va.maximum_scale = 870
va.tick_labels.font.size = Pt(10); va.tick_labels.font.color.rgb = INK_SOFT
va.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xEC, 0xEC)
va.format.line.color.rgb = GREY_L
ca = ch.category_axis
ca.tick_labels.font.size = Pt(10.5); ca.tick_labels.font.color.rgb = INK
ca.major_tick_mark = XL_TICK_MARK.NONE
para_block(s, 0.95, 5.30, 6.90, 0.30,
           [dict(text='평탄 2 → 5 :  +17.8 µm / 1,626 ms  =  0.0110 µm/ms', size=12.5, bold=True, color=NAVY)])
RX4, RW4 = 8.15, 4.63
para_block(s, RX4, 2.08, RW4, 0.30, [dict(text='네 상한에 실측을 대입한 여유', size=13, bold=True, color=ACCENT_D)])
MARG = [('① FASTECH 500k pulse/s', '50.0', '4,560배'), ('② Vel 150,000', '15.0', '1,368배'),
        ('②′ Accel 200 M', '10.0', '912배'), ('③ 초점 정확도', '10.0', '912배')]
tb = table(s, RX4, 2.46, RW4, 2.20, 5, 3, col_w=[2.4, 1.0, 1.2], header_h=0.44, row_h=0.44)
for c, h in enumerate(['상한', '한계\n[µm/ms]', '평균 대비\n여유']):
    cell(tb, 0, c, h, 10.5, True, WHITE, ACCENT_D, PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
for r, row in enumerate(MARG, 1):
    hot = r >= 3
    for c, v in enumerate(row):
        cell(tb, r, c, v, 11, hot, NAVY if hot else INK,
             ACCENT_L if hot else (RGBColor(0xF7, 0xF9, 0xFC) if r % 2 else None),
             PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
style_table(tb)
para_block(s, RX4, 4.78, RW4, 0.50,
           [dict(text='실측 요구 변화율 0.0110 µm/ms (X축 환산 0.137 µm/mm)  ·  순시 p99 1.879 µm/ms 는 센서 노이즈 — 지배 상한 대비 순시 여유 5배',
                 size=11.5, line_pct=125)])
NB = [(0.95, 5.85, '적용 시프트 — 계산 900.66 vs 실측 901.5 ms',
       'Lead Control 이 실제로는 거의 걸리지 않는다. 주기당 지령 변화가 중앙값 0.03 µm(0.3 pulse)에 그쳐 선행 시간이 사실상 0 이다 — 오류가 아니라 예시값과 실동작의 차이다.'),
      (7.08, 5.70, '로깅 ON 조건에서의 비용',
       'FOCUS_DEBUG_LOG_PER_CYCLE_ENABLE = 1 인 Debug 빌드에서 주기 초과 5.84 %(538/9,218) · 최대 주기 19,500 µs — 진단 로그를 켠 채 타이밍을 평가하면 안 된다.')]
for x, w, h, t in NB:
    rrect(s, x, 5.66, w, 1.15, CARD, None, radius=0.05)
    para_block(s, x + 0.24, 5.82, w - 0.48, 0.28, [dict(text=h, size=12.5, bold=True, color=ACCENT_D)])
    para_block(s, x + 0.24, 6.14, w - 0.48, 0.58, [dict(text=t, size=11, line_pct=125)])

# ================================================================ 16. 모션 지령
s = new('모션 지령 결과 — 정반대의 두 실패')
kicker(s, '수락 여부를 가르는 것은 부하나 통신이 아니라 그 주기의 이동 거리 하나다')
MOT = [('거부 4151', 'IN_NONMOTION', '322건 · 69.4 %', '0.380 pulse', '0.087 ms',
        '명령이 도착했을 때 이미 멈춰 있었다 (가속 중으로 읽힘 97.8 %)', '너무 빨라서', RED),
       ('수락', 'Override 성공', '36건 · 7.8 %', '5.592 pulse', '0.334 ms',
        '아직 가속 중이라 진로를 바꿀 여유가 남아 있는 좁은 창', '—', RGBColor(0x1F, 0x6B, 0x4A)),
       ('거부 4168', 'CACULATION', '106건 · 22.8 %', '29.022 pulse', '0.762 ms',
        '이미 계산된 감속 곡선을 남은 거리로 다시 짤 수 없었다 (감속 중 89.6 %)', '너무 커서', ORANGE)]
for i, (code, name, cnt, mv, tt, why, tag, col) in enumerate(MOT):
    x = 0.95 + i * (CW3 + GAP)
    rrect(s, x, 2.12, CW3, 3.30, CARD, INK, 0.75, 0.03)
    hd = rrect(s, x, 2.12, CW3, 0.62, col, None, radius=0.03)
    label_in(hd, code + '   ' + name, 13.5, True, WHITE)
    rows = [('건수', cnt), ('이동량 p50', mv), ('소요 시간', tt)]
    yy = 2.92
    for k, v in rows:
        para_block(s, x + 0.26, yy, CW3 - 0.52, 0.28, [dict(text=k, size=11.5, color=INK_SOFT)])
        para_block(s, x + 0.26, yy, CW3 - 0.52, 0.28, [dict(text=v, size=13, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)])
        hline(s, x + 0.26, yy + 0.32, CW3 - 0.52)
        yy += 0.50
    para_block(s, x + 0.26, 4.46, CW3 - 0.52, 0.68, [dict(text=why, size=11.5, line_pct=125)])
    if tag != '—':
        t = rrect(s, x + 0.26, 5.02, 1.45, 0.34, col, None, radius=0.10)
        label_in(t, tag, 11.5, True, WHITE)
para_block(s, 0.95, 5.62, CW, 1.10, [
    dict(text='AxmMoveStartPos(축 정지 중) 5,224건은 전량 수락 — 전체 5,688건 중 수락 92.5 %', size=13, bold=True, color=ACCENT_D, space_after=7),
    dict(text='초점 오차로는 이어지지 않는다. 지령이 절대위치이므로 거부된 값은 누적되지 않고 다음 수락 지령이 올바른 좌표로 데려간다. '
              '실측 회복 지연 p50 0 ms · p95 2 ms · 최대 8 ms, 그동안 목표가 낡는 양은 p95 1.30 µm(D.O.F. 예산의 37 %)로 최대치에서만 초과한다.',
         size=12.5, line_pct=130)])

# ================================================================ 17. 종합 판정
s = new('종합 판정 — 현 요구 충족')
RES = [('29 %', 'D.O.F. 예산 사용', '오차 1.02 µm / 3.5 µm'),
       ('23배', '추종 한계 향상', '0.43 → 10.0 µm/ms'),
       ('912배', '실측 여유', '0.011 µm/ms vs 10 µm/ms')]
for i, (big, head, sub) in enumerate(RES):
    cx = 2.35 + i * 4.15
    rrect(s, cx - 0.92, 2.00, 1.84, 1.84, ACCENT_L, ACCENT_D, 2.5, shape=MSO_SHAPE.OVAL)
    para_block(s, cx - 1.40, 2.60, 2.80, 0.66,
               [dict(text=big, size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=TITLE_FONT)])
    para_block(s, cx - 1.75, 4.00, 3.50, 0.32, [dict(text=head, size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)])
    para_block(s, cx - 1.75, 4.34, 3.50, 0.30, [dict(text=sub, size=12, color=INK_SOFT, align=PP_ALIGN.CENTER)])
GOAL = [('촬상 시점 초점 오차', '≤ 3.5 µm (카메라 D.O.F.)', '1.01 µm · 예산의 29 %', '달성'),
        ('지연 재현 정확도', '≤ 0.25 ms', '±0.25 ms (반주기 선보정)', '달성'),
        ('제어 주기', '≤ 1 ms', '500 µs · 38,007주기 초과 0회', '달성'),
        ('Z축 추종 속도 · 부하율', '≥ 10 µm/ms · 상한의 20 % 이내', '10.0 µm/ms · 실동작 17 %', '달성')]
tb = table(s, 0.95, 5.02, CW, 1.55, 5, 4, col_w=[3.0, 3.4, 4.0, 1.1], header_h=0.36, row_h=0.34)
for c, h in enumerate(['개발 목표', '기준', '검증 결과', '판정']):
    cell(tb, 0, c, h, 11.5, True, WHITE, ACCENT_D, PP_ALIGN.CENTER if c == 3 else PP_ALIGN.LEFT)
for r, row in enumerate(GOAL, 1):
    for c, v in enumerate(row):
        cell(tb, r, c, v, 11, c == 3, RGBColor(0x1F, 0x6B, 0x4A) if c == 3 else INK,
             RGBColor(0xF7, 0xF9, 0xFC) if r % 2 else None,
             PP_ALIGN.CENTER if c == 3 else PP_ALIGN.LEFT)
style_table(tb)
para_block(s, 0.95, 4.72, CW - 0.10, 0.30,
    [dict(text='추종 한계 10.0 µm/ms 는 실측 조건 0.04 µm/ms 의 250배, 개발 목표 0.08 µm/ms 의 125배 — 이 여유가 「정상 이미지 획득」의 근거다.', size=11.5, color=INK_SOFT)])

# ================================================================ 18. 권고 조치와 확장성
s = new('권고 조치와 확장 가능성')
kicker(s, '현 요구는 이미 충족한다 — 아래는 장애 복구가 아니라 여유 확보가 목적이다')
REC = [(0.95, 5.70, 'A. 측정으로 결정할 설정', '코드 수정이 아니라 검증 테스트로 결정', [
            '정지 상태 전압 노이즈(p-p) 실측 하나로 데드밴드·필터를 동시에 판단',
            '노이즈가 ±1 µm(±5 mV)보다 작으면 데드밴드 SKIP 이 유리 — 예산 1.00 µm 회수',
            '필터는 군지연이 보상되어 정상상태 기여 0 — 근거는 노이즈와 계단 응답',
            'Lead Control 은 유지 — SKIP 시 Δt 1.09 ms, 추종 한계 2.3 µm/ms 로 하락']),
       (7.08, 5.70, 'B. 10 µm/ms 를 사양화할 경우', '여유가 정확히 0 인 균형점이라 둘 중 하나 필요', [
            '① 데드밴드 SKIP — 비용 없이 1 µm 회수해 14 µm/ms (우선 권고)',
            '② 제어 주기 250 µs — 가속도 400 M 동반 필요 (작업시간 211 µs = 84 %)',
            '계단 전이가 문제되면 대칭(zero-phase) 필터 — 추가 지연 비용 없음',
            'Setting.ini [ITEM2] 미반영 — 2번 카메라 운용 시 함께 맞출 것'])]
for x, w, head, sub, lines in REC:
    rrect(s, x, 2.12, w, 3.28, CARD, INK, 0.75, 0.03)
    para_block(s, x + 0.28, 2.32, w - 0.55, 0.32, [dict(text=head, size=16, bold=True, color=NAVY)])
    para_block(s, x + 0.28, 2.70, w - 0.55, 0.28, [dict(text=sub, size=11.5, color=ACCENT_D)])
    para_block(s, x + 0.28, 3.06, w - 0.55, 2.10,
               [dict(text=t, size=12, bullet='·', line_pct=120, space_after=8) for t in lines])
CHIP = [('미확보 — Z축 실위치', '엔코더 되먹임이 없어 AxmStatusGetActPos 는 0.000. Confocal 기준 또는 외부 계측기 필요'),
        ('미확보 — 초기 오프셋 분포', '패널 간 초기 높이 편차. 10장 정도의 초기 거리값 기록으로 확인된다'),
        ('확장 — 스캔 속도 증속', '지연 상수 72.5 mm ÷ v 재계산. 120 mm/s → 604 ms · 추종 사양 동시 갱신')]
for i, (h, t) in enumerate(CHIP):
    x = 0.95 + i * (CW3 + GAP)
    rrect(s, x, 5.54, CW3, 1.00, CARD_B, None, radius=0.05)
    para_block(s, x + 0.22, 5.70, CW3 - 0.44, 0.26, [dict(text=h, size=12, bold=True, color=ACCENT_D)])
    para_block(s, x + 0.22, 6.00, CW3 - 0.44, 0.46, [dict(text=t, size=10.5, line_pct=120)])
footnote(s, '구조가 지연 상수 하나로 닫히므로, 스캔 속도와 선행 거리만 바꾸면 동일 설계를 타 검사 장비에 그대로 이식할 수 있다 — 전사 광학 검사 장비의 표준 초점 모듈로 확장 가능하다.', 6.68)

# ================================================================ 마무리
lst = prs.slides._sldIdLst
el = list(lst)[1]                    # '감사 합니다' 슬라이드를 맨 뒤로
lst.remove(el); lst.append(el)

prs.save(OUT)
print('saved', OUT, '· slides =', len(prs.slides.__iter__.__self__._sldIdLst))
