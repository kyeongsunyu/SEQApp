# -*- coding: utf-8 -*-
"""제1회 연구소기술경진대회 양식(한국전광 OLB 통합검사기 덱) 재현용 헬퍼."""
import copy
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- 양식 토큰 (원본 덱에서 추출) ----------------------------------------
NAVY      = RGBColor(0x00, 0x20, 0x60)   # 제목 글자색
INK       = RGBColor(0x38, 0x38, 0x38)   # 본문 글자색
INK_SOFT  = RGBColor(0x66, 0x66, 0x66)
ACCENT    = RGBColor(0x5B, 0x9B, 0xD5)   # theme accent1
ACCENT_D  = RGBColor(0x44, 0x78, 0xA0)   # accent1 lumMod 75%
ACCENT_L  = RGBColor(0xDE, 0xEA, 0xF6)
CARD      = RGBColor(0xF2, 0xEE, 0xEE)   # 원본 카드 배경
CARD_B    = RGBColor(0xEB, 0xF2, 0xF9)
GREEN     = RGBColor(0x92, 0xD0, 0x50)   # 원본 섹션 라벨 pill
ORANGE    = RGBColor(0xE0, 0x4F, 0x00)
RED       = RGBColor(0xC0, 0x30, 0x2A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY_L    = RGBColor(0xD9, 0xD9, 0xD9)

TITLE_FONT = 'HY견고딕'
BODY_FONT  = '맑은 고딕'

SLIDE_W, SLIDE_H = 13.333, 7.5
CX0, CX1 = 0.85, 12.78          # 콘텐츠 좌/우 경계
CW = CX1 - CX0
CY0, CY1 = 1.92, 6.98           # 콘텐츠 상/하 경계

GRAD_CHEVRON = (
 '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
 ' flip="none" rotWithShape="1"><a:gsLst>'
 '<a:gs pos="0"><a:schemeClr val="accent1"><a:tint val="66000"/><a:satMod val="160000"/></a:schemeClr></a:gs>'
 '<a:gs pos="50000"><a:schemeClr val="accent1"><a:tint val="44500"/><a:satMod val="160000"/></a:schemeClr></a:gs>'
 '<a:gs pos="100000"><a:schemeClr val="accent1"><a:tint val="23500"/><a:satMod val="160000"/></a:schemeClr></a:gs>'
 '</a:gsLst><a:lin ang="0" scaled="1"/><a:tileRect/></a:gradFill>')


# ---- 저수준 유틸 ----------------------------------------------------------
def _swap_fill(shape, xml):
    spPr = shape._element.spPr
    for tag in ('solidFill', 'noFill', 'gradFill', 'blipFill', 'pattFill', 'grpFill'):
        el = spPr.find(qn('a:' + tag))
        if el is not None:
            spPr.remove(el)
    frag = etree.fromstring(xml)
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(frag)
    else:
        spPr.append(frag)


def no_line(shape):
    shape.line.fill.background()


def txt_width(text, pt_size, ko=1.00, en=0.52):
    """대략적인 텍스트 폭(inch). 한글은 전각, 라틴은 반각 가정."""
    em = pt_size / 72.0
    w = 0.0
    for ch in text:
        w += em * (ko if ord(ch) > 0x2000 else en)
    return w


def set_para(p, size, bold=False, color=INK, align=PP_ALIGN.LEFT,
             line_pct=None, line_pts=None, space_after=0, bullet=None,
             font=BODY_FONT, italic=False):
    p.alignment = align
    pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    if line_pct:
        ln = etree.SubElement(pPr, qn('a:lnSpc'))
        etree.SubElement(ln, qn('a:spcPct')).set('val', str(int(line_pct * 1000)))
    elif line_pts:
        ln = etree.SubElement(pPr, qn('a:lnSpc'))
        etree.SubElement(ln, qn('a:spcPts')).set('val', str(int(line_pts * 100)))
    if space_after:
        sa = etree.SubElement(pPr, qn('a:spcAft'))
        etree.SubElement(sa, qn('a:spcPts')).set('val', str(int(space_after * 100)))
    if bullet:
        marL = int(Pt(size * 1.15))
        pPr.set('marL', str(marL)); pPr.set('indent', str(-marL))
        etree.SubElement(pPr, qn('a:buSzPct')).set('val', '90000')
        etree.SubElement(pPr, qn('a:buChar')).set('char', bullet)
    for r in p.runs:
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color
        f.name = font
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font)


def textbox(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def para_block(slide, x, y, w, h, items, wrap=True, anchor=MSO_ANCHOR.TOP):
    """items: list of dict(text=, size=, bold=, color=, bullet=, line_pct=,
                            space_after=, align=, font=, italic=)"""
    tb, tf = textbox(slide, x, y, w, h, wrap, anchor)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.add_run().text = it['text']
        set_para(p,
                 it.get('size', 14), it.get('bold', False), it.get('color', INK),
                 it.get('align', PP_ALIGN.LEFT), it.get('line_pct'), it.get('line_pts'),
                 it.get('space_after', 0), it.get('bullet'),
                 it.get('font', BODY_FONT), it.get('italic', False))
    return tb


def rrect(slide, x, y, w, h, fill=CARD, line=None, lw=0.75, radius=0.02,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        no_line(sh)
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.word_wrap = True
    return sh


def rect(slide, x, y, w, h, fill=CARD, line=None, lw=0.75):
    return rrect(slide, x, y, w, h, fill, line, lw, shape=MSO_SHAPE.RECTANGLE)


def pill(slide, x, y, w, h, text, fill=GREEN, size=14, color=RGBColor(0, 0, 0)):
    sh = rrect(slide, x, y, w, h, fill, None,
               shape=MSO_SHAPE.FLOWCHART_ALTERNATE_PROCESS)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.add_run().text = text
    set_para(p, size, True, color, PP_ALIGN.CENTER)
    return sh


def label_in(shape, text, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             font=BODY_FONT, line_pct=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    first = True
    for line in text.split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.add_run().text = line
        set_para(p, size, bold, color, align, line_pct=line_pct, font=font)
    return shape


def arrow(slide, x, y, w, h, fill=ACCENT_D, shape=MSO_SHAPE.RIGHT_ARROW):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    no_line(sh); sh.shadow.inherit = False
    return sh


def hline(slide, x, y, w, color=GREY_L, thick=0.022):
    return rect(slide, x, y, w, thick, fill=color)


# ---- 슬라이드 뼈대 --------------------------------------------------------
def title_bar(slide, text, size=26, x=0.79, y=0.60, h=0.92):
    tw = txt_width(text, size)
    w = min(0.62 + tw + 0.86, 11.6)
    ch = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    _swap_fill(ch, GRAD_CHEVRON)
    no_line(ch); ch.shadow.inherit = False
    ch._element.spPr.find(qn('a:xfrm')).set('flipH', '1')
    tb, tf = textbox(slide, x + 0.55, y, tw + 0.5, h, wrap=False,
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.add_run().text = text
    set_para(p, size, False, NAVY, PP_ALIGN.LEFT, font=TITLE_FONT)
    return ch


def kicker(slide, text, y=1.60, x=0.95, size=12.5, color=INK_SOFT):
    """제목 아래 한 줄 리드 문장."""
    return para_block(slide, x, y, CW, 0.34,
                      [dict(text=text, size=size, color=color, line_pct=100)])


def footnote(slide, text, y=None, size=10.5, color=INK_SOFT, align=PP_ALIGN.LEFT):
    y = CY1 - 0.30 if y is None else y
    return para_block(slide, 0.95, y, CW - 0.10, 0.30,
                      [dict(text=text, size=size, color=color, align=align,
                            line_pct=110)])


def banner(slide, y, text, size=12.5, fill=ACCENT_L, h=0.46, color=NAVY,
           x=CX0, w=None):
    w = CW if w is None else w
    sh = rrect(slide, x, y, w, h, fill, None, radius=0.10)
    label_in(sh, text, size, True, color, PP_ALIGN.CENTER)
    return sh


# ---- 표 -------------------------------------------------------------------
NO_STYLE = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'


def table(slide, x, y, w, h, rows, cols, col_w=None, header_h=0.42, row_h=None):
    gf = slide.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    tbl = gf.table
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '1'); tblPr.set('bandRow', '0')
    st = tblPr.find(qn('a:tableStyleId'))
    if st is None:
        st = etree.SubElement(tblPr, qn('a:tableStyleId'))
    st.text = NO_STYLE
    if col_w:
        tot = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / tot))
    tbl.rows[0].height = Inches(header_h)
    if row_h:
        for r in range(1, rows):
            tbl.rows[r].height = Inches(row_h)
    return tbl


def cell(tbl, r, c, text, size=11.5, bold=False, color=INK, fill=None,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=BODY_FONT,
         line_pct=105):
    cl = tbl.cell(r, c)
    cl.margin_left = cl.margin_right = Inches(0.08)
    cl.margin_top = cl.margin_bottom = Inches(0.03)
    cl.vertical_anchor = anchor
    if fill is None:
        cl.fill.background()
    else:
        cl.fill.solid(); cl.fill.fore_color.rgb = fill
    tf = cl.text_frame
    tf.word_wrap = True
    first = True
    for line in str(text).split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.add_run().text = line
        set_para(p, size, bold, color, align, line_pct=line_pct, font=font)
    return cl


def cell_border(tbl, r, c, edges='B', color=GREY_L, w=0.75):
    """edges: 문자열 조합 'LRTB'"""
    tc = tbl.cell(r, c)._tc
    tcPr = tc.get_or_add_tcPr()
    order = {'L': 'a:lnL', 'R': 'a:lnR', 'T': 'a:lnT', 'B': 'a:lnB'}
    for e in 'LRTB':
        if e not in edges:
            continue
        tag = order[e]
        old = tcPr.find(qn(tag))
        if old is not None:
            tcPr.remove(old)
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set('w', str(int(Pt(w)))); ln.set('cap', 'flat')
        ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr')).set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
    # OOXML 스키마 순서(lnL,lnR,lnT,lnB,...)로 재정렬
    for tag in ('a:lnB', 'a:lnT', 'a:lnR', 'a:lnL'):
        el = tcPr.find(qn(tag))
        if el is not None:
            tcPr.remove(el); tcPr.insert(0, el)


def style_table(tbl, head_fill=ACCENT_D, head_color=WHITE, zebra=RGBColor(0xF6, 0xF8, 0xFB)):
    nrows = len(tbl.rows); ncols = len(tbl.columns)
    for c in range(ncols):
        cell_border(tbl, 0, c, 'B', RGBColor(0x2C, 0x31, 0x4F), 1.2)
    for r in range(1, nrows):
        for c in range(ncols):
            cell_border(tbl, r, c, 'B', GREY_L, 0.6)
